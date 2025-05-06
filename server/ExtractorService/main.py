from fastapi import FastAPI, Body
from typing import List, Dict, Any
import boto3, tempfile, os, logging
from youtube_transcript_api import YouTubeTranscriptApi
from pptx import Presentation
from PyPDF2 import PdfReader
import nltk
import bisect


# Ensure punkt tokenizer is available
nltk.download('punkt')
from nltk.tokenize import sent_tokenize

# ========== Logging Setup ==========
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s | %(levelname)s | %(message)s",
)
logger = logging.getLogger(__name__)

# ========== FastAPI App ==========
app = FastAPI()
s3 = boto3.client('s3')
from fastapi.middleware.cors import CORSMiddleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Allows all origins
    allow_credentials=True,
    allow_methods=["*"],  # Allows all methods (GET, POST, PUT, DELETE, etc.)
    allow_headers=["*"],  # Allows all headers
)
# ========== Helpers ==========

def chunk_text(text: str, max_tokens=300, overlap=50) -> List[str]:
    sentences = sent_tokenize(text)
    chunks = []
    chunk = []
    token_count = 0
    i = 0
    while i < len(sentences):
        sent = sentences[i]
        sent_tokens = len(sent.split())
        if token_count + sent_tokens <= max_tokens:
            chunk.append(sent)
            token_count += sent_tokens
            i += 1
        else:
            chunks.append(" ".join(chunk))
            # overlap logic
            overlap_chunk = []
            overlap_token_count = 0
            j = len(chunk) - 1
            while j >= 0 and overlap_token_count < overlap:
                overlap_chunk.insert(0, chunk[j])
                overlap_token_count += len(chunk[j].split())
                j -= 1
            chunk = overlap_chunk
            token_count = overlap_token_count
    if chunk:
        chunks.append(" ".join(chunk))
    return chunks


def parse_pdf(file_path: str) -> str:
    logger.info(f"Parsing PDF: {file_path}")
    reader = PdfReader(file_path)
    return "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])

def parse_pptx(file_path: str) -> str:
    logger.info(f"Parsing PPTX: {file_path}")
    prs = Presentation(file_path)
    return "\n".join([
        shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")
    ])

# ========== Main API ==========

@app.post("/process")
async def process_docs(payload: Dict[str, List[str]] = Body(...)) -> Dict[str, Any]:
    documents = payload.get("documents", [])
    videos = payload.get("youtube_videos", [])
    results = {"document_chunks": [], "youtube_chunks": [], "errors": []}

    # S3 documents
    for s3_uri in documents:
        try:
            logger.info(f"Processing S3 document: {s3_uri}")
            if not s3_uri.startswith("s3://"):
                raise ValueError(f"Invalid S3 URI: {s3_uri}")
            bucket, key = s3_uri[5:].split("/", 1)
            ext = os.path.splitext(key)[1][1:].lower()

            with tempfile.NamedTemporaryFile(delete=False) as tmp:
                s3.download_fileobj(bucket, key, tmp)
                tmp.flush()
                logger.info(f"Downloaded S3 file: {s3_uri}")

                if ext == "pdf":
                    text = parse_pdf(tmp.name)
                elif ext == "pptx":
                    text = parse_pptx(tmp.name)
                elif ext == "txt":
                    with open(tmp.name, 'r') as f:
                        text = f.read()
                    logger.info(f"Read plain text file: {s3_uri}")
                else:
                    raise ValueError(f"Unsupported file type: {ext}")

                chunks = chunk_text(text)
                results["document_chunks"].append({
                    "s3_uri": s3_uri,
                    "chunks": chunks
                })
                logger.info(f"Document chunking complete for: {s3_uri}")
        except Exception as e:
            logger.error(f"Error processing document {s3_uri}: {e}")
            results["errors"].append({"document": s3_uri, "error": str(e)})
        finally:
            if os.path.exists(tmp.name):
                os.unlink(tmp.name)
                logger.info(f"Temporary file deleted: {tmp.name}")

    # YouTube videos
    for vid in videos:
        try:
            logger.info(f"Fetching transcript for YouTube video: {vid}")
            transcript = YouTubeTranscriptApi.get_transcript(vid)
            full_text = " ".join([seg["text"] for seg in transcript])
            chunks = chunk_text(full_text)
            results["youtube_chunks"].append({
                "video_id": vid,
                "chunks": chunks
            })
            logger.info(f"YouTube transcript processed: {vid}")
        except Exception as e:
            logger.error(f"Error processing video {vid}: {e}")
            results["errors"].append({"video_id": vid, "error": str(e)})

    return results

@app.get("/health")
async def health_check() -> Dict[str, str]:
    logger.info("Health check requested")
    return {"status": "ok"}

@app.get("/")
async def root() -> Dict[str, str]:
    logger.info("Root endpoint hit")
    return {"message": "Welcome to the Extractor Service!"}

if __name__ == "__main__":
    import uvicorn
    logger.info("Starting Extractor Service.....")
    uvicorn.run(app, host="0.0.0.0", port=8001)
