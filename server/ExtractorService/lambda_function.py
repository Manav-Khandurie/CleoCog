import boto3
import os
import tempfile
from typing import Dict, List, Union
from youtube_transcript_api import YouTubeTranscriptApi
from pptx import Presentation
from PyPDF2 import PdfReader

s3 = boto3.client('s3')

def lambda_handler(event: Dict, context: object) -> Dict:
    """
    Process documents from S3 and YouTube videos for transcription.
    
    Expected event format:
    {
        "documents": ["s3://bucket-name/path/to/file.pdf", ...],
        "youtube_videos": ["VIDEO_ID_1", "VIDEO_ID_2", ...]
    }
    """
    results = {
        "document_texts": [],
        "youtube_transcripts": [],
        "errors": []
    }

    # Process documents
    for s3_uri in event.get('documents', []):
        try:
            # Extract bucket and key from S3 URI
            if not s3_uri.startswith('s3://'):
                raise ValueError(f"Invalid S3 URI: {s3_uri}")
            bucket, key = s3_uri[5:].split('/', 1)
            
            # Download file
            ext = os.path.splitext(key)[1][1:].lower()
            with tempfile.NamedTemporaryFile(delete=False) as tmp_file:
                s3.download_fileobj(bucket, key, tmp_file)
                tmp_file.flush()
                
                # Process based on file type
                if ext == 'pdf':
                    text = parse_pdf(tmp_file.name)
                elif ext == 'pptx':
                    text = parse_pptx(tmp_file.name)
                elif ext == 'txt':
                    with open(tmp_file.name, 'r') as f:
                        text = f.read()
                else:
                    raise ValueError(f"Unsupported file type: {ext}")
                
                results["document_texts"].append({
                    "s3_uri": s3_uri,
                    "content": text
                })
                
        except Exception as e:
            results["errors"].append({
                "document": s3_uri,
                "error": str(e)
            })
        finally:
            if tmp_file:
                os.unlink(tmp_file.name)

    # Process YouTube videos
    for video_id in event.get('youtube_videos', []):
        try:
            transcript = YouTubeTranscriptApi.get_transcript(video_id)
            results["youtube_transcripts"].append({
                "video_id": video_id,
                "transcript": transcript
            })
        except Exception as e:
            results["errors"].append({
                "video_id": video_id,
                "error": str(e)
            })

    return results

def parse_pdf(file_path: str) -> str:
    """Extract text from PDF file"""
    text = []
    reader = PdfReader(file_path)
    for page in reader.pages:
        text.append(page.extract_text())
    return "\n".join(text)

def parse_pptx(file_path: str) -> str:
    """Extract text from PowerPoint file"""
    text = []
    prs = Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)